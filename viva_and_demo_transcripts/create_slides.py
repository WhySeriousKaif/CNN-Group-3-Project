import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    # Initialize Presentation
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette Constants
    BG_COLOR = RGBColor(11, 15, 25)         # Deep Navy/Slate
    CARD_BG = RGBColor(20, 27, 45)          # Dark Card Fill
    CARD_BORDER = RGBColor(180, 115, 20)     # Bronze/Gold Border
    BORDER_SUBTLE = RGBColor(40, 50, 75)    # Subtle Slate Border
    TEXT_WHITE = RGBColor(255, 255, 255)    # Primary Text
    TEXT_GRAY = RGBColor(209, 213, 219)     # Body Text
    TEXT_MUTED = RGBColor(148, 163, 184)    # Subtitle/Cap Text
    ACCENT_GOLD = RGBColor(251, 191, 36)    # Warm Gold Highlight
    TABLE_HEADER_BG = RGBColor(180, 115, 20)

    # Fonts
    HEADER_FONT = "Georgia"
    BODY_FONT = "Calibri"

    def apply_background(slide):
        """Applies solid dark background to the slide."""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def add_slide_header(slide, title_text):
        """Adds a standard slide header with a gold separator line."""
        apply_background(slide)
        
        # Title box
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = HEADER_FONT
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = ACCENT_GOLD
        
        # Horizontal Separator line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.3), Inches(11.733), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = CARD_BORDER
        line.line.fill.background()

    def add_card(slide, left, top, width, height, border_color=BORDER_SUBTLE):
        """Draws a rounded rectangle background card."""
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)
        return card

    def style_paragraph(p, text, font_name, size_pt, color_rgb, bold=False, italic=False, space_after=6):
        """Helper to style paragraph text."""
        p.text = text
        p.font.name = font_name
        p.font.size = Pt(size_pt)
        p.font.color.rgb = color_rgb
        p.font.bold = bold
        p.font.italic = italic
        p.space_after = Pt(space_after)
        p.line_spacing = 1.15

    # ==========================================
    # SLIDE 1: Title Slide
    # ==========================================
    slide_layout = prs.slide_layouts[6] # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    apply_background(slide)

    # Title box
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.733), Inches(1.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    style_paragraph(p, "Optimization Dynamics and Sequence Generation in Deep Learning", HEADER_FONT, 38, ACCENT_GOLD, bold=True)
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    p_sub = tf.add_paragraph()
    style_paragraph(p_sub, "Comparative Analysis of Batch Normalization & Character-Level LSTMs", BODY_FONT, 20, TEXT_MUTED, italic=True)
    p_sub.alignment = PP_ALIGN.CENTER

    # Divider Line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(4.5), Inches(10.333), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = CARD_BORDER
    line.line.fill.background()

    # Column 1 Team members
    team1_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.8), Inches(5.0), Inches(2.2))
    tf1 = team1_box.text_frame
    tf1.word_wrap = True
    
    p1 = tf1.paragraphs[0]
    style_paragraph(p1, "MD Kaif Molla", HEADER_FONT, 14, TEXT_WHITE, bold=True, space_after=2)
    p1_sub = tf1.add_paragraph()
    style_paragraph(p1_sub, "md.24bcs10221@sst.scaler.com", BODY_FONT, 12, TEXT_MUTED, space_after=12)
    
    p2 = tf1.add_paragraph()
    style_paragraph(p2, "Arhan Das", HEADER_FONT, 14, TEXT_WHITE, bold=True, space_after=2)
    p2_sub = tf1.add_paragraph()
    style_paragraph(p2_sub, "Arhan.24bcs10023@sst.scaler.com", BODY_FONT, 12, TEXT_MUTED, space_after=12)

    p3 = tf1.add_paragraph()
    style_paragraph(p3, "Sumit Akhuli", HEADER_FONT, 14, TEXT_WHITE, bold=True, space_after=2)
    p3_sub = tf1.add_paragraph()
    style_paragraph(p3_sub, "Sumit.24bcs10158@sst.scaler.com", BODY_FONT, 12, TEXT_MUTED)

    # Column 2 Team members
    team2_box = slide.shapes.add_textbox(Inches(6.8), Inches(4.8), Inches(5.0), Inches(2.2))
    tf2 = team2_box.text_frame
    tf2.word_wrap = True
    
    p4 = tf2.paragraphs[0]
    style_paragraph(p4, "Sarthak Mishra", HEADER_FONT, 14, TEXT_WHITE, bold=True, space_after=2)
    p4_sub = tf2.add_paragraph()
    style_paragraph(p4_sub, "Sarthak.24bcs10470@sst.scaler.com", BODY_FONT, 12, TEXT_MUTED, space_after=12)
    
    p5 = tf2.add_paragraph()
    style_paragraph(p5, "Netram Faran", HEADER_FONT, 14, TEXT_WHITE, bold=True, space_after=2)
    p5_sub = tf2.add_paragraph()
    style_paragraph(p5_sub, "netram.24bcs10329@sst.scaler.com", BODY_FONT, 12, TEXT_MUTED, space_after=12)

    p6 = tf2.add_paragraph()
    style_paragraph(p6, "Aditya Prasad", HEADER_FONT, 14, TEXT_WHITE, bold=True, space_after=2)
    p6_sub = tf2.add_paragraph()
    style_paragraph(p6_sub, "aditya.24bcs10179@sst.scaler.com", BODY_FONT, 12, TEXT_MUTED)

    # ==========================================
    # SLIDE 2: Task 1: AI Story Continuation Overview
    # ==========================================
    slide = prs.slides.add_slide(slide_layout)
    add_slide_header(slide, "Task 1: AI Story Continuation Overview")

    # Card 1: Objective
    add_card(slide, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.0), CARD_BORDER)
    box1 = slide.shapes.add_textbox(Inches(1.05), Inches(1.85), Inches(5.1), Inches(4.5))
    tf = box1.text_frame
    tf.word_wrap = True
    style_paragraph(tf.paragraphs[0], "Project Objective", HEADER_FONT, 18, ACCENT_GOLD, bold=True, space_after=12)
    
    p = tf.add_paragraph()
    style_paragraph(p, "• Train a character-level Stacked LSTM network to generate story continuations from text prompts.", BODY_FONT, 14, TEXT_GRAY, space_after=12)
    
    p = tf.add_paragraph()
    style_paragraph(p, "• Forces the model to learn spelling, grammar, punctuation, and style from scratch, without using pre-defined word tokens.", BODY_FONT, 14, TEXT_GRAY, space_after=12)
    
    p = tf.add_paragraph()
    style_paragraph(p, "• Serves as an illustration of sequence representation learning directly from raw text sequences.", BODY_FONT, 14, TEXT_GRAY)

    # Card 2: Core Questions Explored
    add_card(slide, Inches(6.933), Inches(1.6), Inches(5.6), Inches(5.0), BORDER_SUBTLE)
    box2 = slide.shapes.add_textbox(Inches(7.183), Inches(1.85), Inches(5.1), Inches(4.5))
    tf = box2.text_frame
    tf.word_wrap = True
    style_paragraph(tf.paragraphs[0], "Core Questions Explored", HEADER_FONT, 18, ACCENT_GOLD, bold=True, space_after=12)

    p = tf.add_paragraph()
    style_paragraph(p, "• Representation Learning: Can a model learn English syntax and style entirely at the character level?", BODY_FONT, 14, TEXT_GRAY, space_after=12)

    p = tf.add_paragraph()
    style_paragraph(p, "• Decoding Strategies: How does deterministic greedy decoding compare to temperature-based probabilistic sampling?", BODY_FONT, 14, TEXT_GRAY, space_after=12)

    p = tf.add_paragraph()
    style_paragraph(p, "• Stylistic Nuance: How does the style of the training corpus affect the vocabulary, tone, and pacing of the generated output?", BODY_FONT, 14, TEXT_GRAY)


    # ==========================================
    # SLIDE 3: Sequence Modeling: LSTM Architecture
    # ==========================================
    slide = prs.slides.add_slide(slide_layout)
    add_slide_header(slide, "Sequence Modeling: LSTM Architecture")

    # Card 1: Why LSTM & Gating
    add_card(slide, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.0), BORDER_SUBTLE)
    box1 = slide.shapes.add_textbox(Inches(1.05), Inches(1.85), Inches(5.1), Inches(4.5))
    tf = box1.text_frame
    tf.word_wrap = True
    style_paragraph(tf.paragraphs[0], "Why LSTM & Gating Mechanism?", HEADER_FONT, 18, ACCENT_GOLD, bold=True, space_after=8)

    p = tf.add_paragraph()
    style_paragraph(p, "• Vanilla RNNs suffer from vanishing and exploding gradients due to repeated temporal multiplications.", BODY_FONT, 12, TEXT_GRAY, space_after=8)

    p = tf.add_paragraph()
    style_paragraph(p, "• LSTM gates control information flow through cell state (Cₜ):", BODY_FONT, 12, TEXT_GRAY, space_after=8)

    p = tf.add_paragraph()
    style_paragraph(p, "  - Forget Gate: Decides what history to discard.", BODY_FONT, 12, TEXT_WHITE, bold=True, space_after=2)
    p = tf.add_paragraph()
    style_paragraph(p, "      fₜ = σ(W_f · [hₜ₋₁, xₜ] + b_f)", HEADER_FONT, 13, ACCENT_GOLD, bold=True, space_after=6)

    p = tf.add_paragraph()
    style_paragraph(p, "  - Input Gate: Decides what new details to store in cell state.", BODY_FONT, 12, TEXT_WHITE, bold=True, space_after=2)
    p = tf.add_paragraph()
    style_paragraph(p, "      iₜ = σ(W_i · [hₜ₋₁, xₜ] + b_i)", HEADER_FONT, 13, ACCENT_GOLD, bold=True, space_after=6)

    p = tf.add_paragraph()
    style_paragraph(p, "  - Output Gate: Decides what hidden state (hₜ) to output.", BODY_FONT, 12, TEXT_WHITE, bold=True, space_after=2)
    p = tf.add_paragraph()
    style_paragraph(p, "      oₜ = σ(W_o · [h₁₋₁, xₜ] + b_o)", HEADER_FONT, 13, ACCENT_GOLD, bold=True)

    # Card 2: Network Configuration
    add_card(slide, Inches(6.933), Inches(1.6), Inches(5.6), Inches(5.0), CARD_BORDER)
    box2 = slide.shapes.add_textbox(Inches(7.183), Inches(1.85), Inches(5.1), Inches(4.5))
    tf = box2.text_frame
    tf.word_wrap = True
    style_paragraph(tf.paragraphs[0], "Our LSTM Network Configuration", HEADER_FONT, 18, ACCENT_GOLD, bold=True, space_after=12)

    p = tf.add_paragraph()
    style_paragraph(p, "• Embedding Layer:", BODY_FONT, 14, TEXT_WHITE, bold=True, space_after=2)
    p = tf.add_paragraph()
    style_paragraph(p, "  - Projects character indices to 128-dimensional dense vectors.", BODY_FONT, 13, TEXT_GRAY, space_after=12)

    p = tf.add_paragraph()
    style_paragraph(p, "• Stacked LSTM Layers:", BODY_FONT, 14, TEXT_WHITE, bold=True, space_after=2)
    p = tf.add_paragraph()
    style_paragraph(p, "  - 2 layers with 256 hidden units each.\n  - Dropout rate of 0.2 after each layer to prevent overfitting.", BODY_FONT, 13, TEXT_GRAY, space_after=12)

    p = tf.add_paragraph()
    style_paragraph(p, "• Linear Decoder Layer:", BODY_FONT, 14, TEXT_WHITE, bold=True, space_after=2)
    p = tf.add_paragraph()
    style_paragraph(p, "  - Projects the 256-dim hidden states to character vocabulary logits.", BODY_FONT, 13, TEXT_GRAY)


    # ==========================================
    # SLIDE 4: Data Pipeline & Corpuses
    # ==========================================
    slide = prs.slides.add_slide(slide_layout)
    add_slide_header(slide, "Data Pipeline & Corpuses")

    # Card 1: Corpuses
    add_card(slide, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.0), BORDER_SUBTLE)
    box1 = slide.shapes.add_textbox(Inches(1.05), Inches(1.85), Inches(5.1), Inches(4.5))
    tf = box1.text_frame
    tf.word_wrap = True
    style_paragraph(tf.paragraphs[0], "Contrasting Training Corpuses", HEADER_FONT, 18, ACCENT_GOLD, bold=True, space_after=12)

    p = tf.add_paragraph()
    style_paragraph(p, "• Shakespeare Corpus (tinyshakespeare):", BODY_FONT, 14, TEXT_WHITE, bold=True, space_after=2)
    p = tf.add_paragraph()
    style_paragraph(p, "  - Poetic syntax, old English vocabulary, dialogue structure.\n  - Vocabulary size: 37 unique characters.", BODY_FONT, 13, TEXT_GRAY, space_after=16)

    p = tf.add_paragraph()
    style_paragraph(p, "• Sherlock Holmes Corpus (Conan Doyle):", BODY_FONT, 14, TEXT_WHITE, bold=True, space_after=2)
    p = tf.add_paragraph()
    style_paragraph(p, "  - 19th-century detective prose, narrative syntax, descriptive style.\n  - Vocabulary size: 57 unique characters.", BODY_FONT, 13, TEXT_GRAY)

    # Card 2: Preprocessing
    add_card(slide, Inches(6.933), Inches(1.6), Inches(5.6), Inches(5.0), CARD_BORDER)
    box2 = slide.shapes.add_textbox(Inches(7.183), Inches(1.85), Inches(5.1), Inches(4.5))
    tf = box2.text_frame
    tf.word_wrap = True
    style_paragraph(tf.paragraphs[0], "Data Preprocessing & Preparation", HEADER_FONT, 18, ACCENT_GOLD, bold=True, space_after=12)

    p = tf.add_paragraph()
    style_paragraph(p, "• Lowercase Normalization:", BODY_FONT, 14, TEXT_WHITE, bold=True, space_after=2)
    p = tf.add_paragraph()
    style_paragraph(p, "  - Reduced all characters to lowercase to keep the vocabulary sizes compact, leading to faster training.", BODY_FONT, 13, TEXT_GRAY, space_after=16)

    p = tf.add_paragraph()
    style_paragraph(p, "• Input-Target Slicing:", BODY_FONT, 14, TEXT_WHITE, bold=True, space_after=2)
    p = tf.add_paragraph()
    style_paragraph(p, "  - Sliced raw text into overlapping sequences of length T = 100.\n  - Target sequence is shifted by one character.\n  - Model learns to predict next character at each time step.", BODY_FONT, 13, TEXT_GRAY)


    # ==========================================
    # SLIDE 5: Decoding and Sampling Techniques
    # ==========================================
    slide = prs.slides.add_slide(slide_layout)
    add_slide_header(slide, "Decoding and Sampling Techniques")

    # Card 1: Greedy vs Temp
    add_card(slide, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.0), CARD_BORDER)
    box1 = slide.shapes.add_textbox(Inches(1.05), Inches(1.85), Inches(5.1), Inches(4.5))
    tf = box1.text_frame
    tf.word_wrap = True
    style_paragraph(tf.paragraphs[0], "Greedy vs. Temperature Sampling", HEADER_FONT, 18, ACCENT_GOLD, bold=True, space_after=12)

    p = tf.add_paragraph()
    style_paragraph(p, "• Greedy Decoding (T = 0):", BODY_FONT, 14, TEXT_WHITE, bold=True, space_after=2)
    p = tf.add_paragraph()
    style_paragraph(p, "  - Always selects the character with the maximum probability.\n  - Highly deterministic, but prone to repeating words and getting stuck in infinite loops.", BODY_FONT, 13, TEXT_GRAY, space_after=12)

    p = tf.add_paragraph()
    style_paragraph(p, "• Temperature Sampling (T > 0):", BODY_FONT, 14, TEXT_WHITE, bold=True, space_after=2)
    p = tf.add_paragraph()
    style_paragraph(p, "  - Scales logits before the Softmax function to adjust entropy:", BODY_FONT, 13, TEXT_GRAY, space_after=4)
    p = tf.add_paragraph()
    style_paragraph(p, "      P(cᵢ) = exp(logits_i / T) / Σⱼ exp(logits_j / T)", HEADER_FONT, 16, ACCENT_GOLD, bold=True)

    # Card 2: Temp settings
    add_card(slide, Inches(6.933), Inches(1.6), Inches(5.6), Inches(5.0), BORDER_SUBTLE)
    box2 = slide.shapes.add_textbox(Inches(7.183), Inches(1.85), Inches(5.1), Inches(4.5))
    tf = box2.text_frame
    tf.word_wrap = True
    style_paragraph(tf.paragraphs[0], "Effects of Temperature Settings", HEADER_FONT, 18, ACCENT_GOLD, bold=True, space_after=12)

    p = tf.add_paragraph()
    style_paragraph(p, "• Low Temperature (T = 0.2):", BODY_FONT, 14, TEXT_WHITE, bold=True, space_after=2)
    p = tf.add_paragraph()
    style_paragraph(p, "  - Very conservative output, highly repetitive, lacks creativity.", BODY_FONT, 13, TEXT_GRAY, space_after=16)

    p = tf.add_paragraph()
    style_paragraph(p, "• Moderate Temperature (T = 0.6):", BODY_FONT, 14, TEXT_WHITE, bold=True, space_after=2)
    p = tf.add_paragraph()
    style_paragraph(p, "  - Optimal balance: creative vocabulary, fluent syntax, and very few spelling/grammar mistakes.", BODY_FONT, 13, TEXT_GRAY, space_after=16)

    p = tf.add_paragraph()
    style_paragraph(p, "• High Temperature (T >= 1.0):", BODY_FONT, 14, TEXT_WHITE, bold=True, space_after=2)
    p = tf.add_paragraph()
    style_paragraph(p, "  - Extreme entropy: highly chaotic, frequent spelling mistakes, outputs look like random gibberish.", BODY_FONT, 13, TEXT_GRAY)


    # ==========================================
    # SLIDE 6: Task 1: Training Loss & Text Samples
    # ==========================================
    slide = prs.slides.add_slide(slide_layout)
    add_slide_header(slide, "Task 1: Training Loss & Text Samples")

    # Left Card: Text info
    add_card(slide, Inches(0.8), Inches(1.6), Inches(5.2), Inches(5.0), BORDER_SUBTLE)
    box1 = slide.shapes.add_textbox(Inches(1.05), Inches(1.85), Inches(4.7), Inches(4.5))
    tf = box1.text_frame
    tf.word_wrap = True
    style_paragraph(tf.paragraphs[0], "Empirical Findings & Text Samples", HEADER_FONT, 18, ACCENT_GOLD, bold=True, space_after=8)

    p = tf.add_paragraph()
    style_paragraph(p, "• Loss Convergence:", BODY_FONT, 13, TEXT_WHITE, bold=True, space_after=2)
    p = tf.add_paragraph()
    style_paragraph(p, "  - Shakespeare Model: final loss of 0.9955 (8 epochs).\n  - Sherlock Model: final loss of 1.1377.", BODY_FONT, 12, TEXT_GRAY, space_after=8)

    p = tf.add_paragraph()
    style_paragraph(p, "• Stylistic Comparison (Prompt: \"the secret of \")", BODY_FONT, 13, TEXT_WHITE, bold=True, space_after=2)
    
    p = tf.add_paragraph()
    style_paragraph(p, "  - Shakespeare (T=0.6):", BODY_FONT, 12, ACCENT_GOLD, bold=True, space_after=2)
    p = tf.add_paragraph()
    style_paragraph(p, "    \"the secret of this too, well of parish that's worth's voice...\"", BODY_FONT, 11, TEXT_GRAY, italic=True, space_after=6)
    
    p = tf.add_paragraph()
    style_paragraph(p, "  - Sherlock Holmes (T=0.6):", BODY_FONT, 12, ACCENT_GOLD, bold=True, space_after=2)
    p = tf.add_paragraph()
    style_paragraph(p, "    \"the secret of home upon his waiting. he woild dost the band...\"", BODY_FONT, 11, TEXT_GRAY, italic=True)

    # Right Card: Image Frame
    add_card(slide, Inches(6.4), Inches(1.6), Inches(6.133), Inches(5.0), CARD_BORDER)
    # Add Image
    img_path = "extracted_plots/task1_lstm_loss.png"
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(6.6), Inches(1.8), Inches(5.733), Inches(4.6))


    # ==========================================
    # SLIDE 7: Task 2: Batch Norm vs. No Batch Norm
    # ==========================================
    slide = prs.slides.add_slide(slide_layout)
    add_slide_header(slide, "Task 2: Batch Norm vs. No Batch Norm")

    # Card 1: Objective
    add_card(slide, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.0), CARD_BORDER)
    box1 = slide.shapes.add_textbox(Inches(1.05), Inches(1.85), Inches(5.1), Inches(4.5))
    tf = box1.text_frame
    tf.word_wrap = True
    style_paragraph(tf.paragraphs[0], "Project Objective", HEADER_FONT, 18, ACCENT_GOLD, bold=True, space_after=12)

    p = tf.add_paragraph()
    style_paragraph(p, "• Train and compare deep feedforward neural networks (6 layers) with and without Batch Normalization on the MNIST dataset.", BODY_FONT, 14, TEXT_GRAY, space_after=12)

    p = tf.add_paragraph()
    style_paragraph(p, "• Systematically evaluate convergence speed, final test accuracy, and stability under different optimization stress-tests.", BODY_FONT, 14, TEXT_GRAY, space_after=12)

    p = tf.add_paragraph()
    style_paragraph(p, "• Provide empirical evidence demonstrating how normalization rescues deep models from training failure.", BODY_FONT, 14, TEXT_GRAY)

    # Card 2: Key Challenges
    add_card(slide, Inches(6.933), Inches(1.6), Inches(5.6), Inches(5.0), BORDER_SUBTLE)
    box2 = slide.shapes.add_textbox(Inches(7.183), Inches(1.85), Inches(5.1), Inches(4.5))
    tf = box2.text_frame
    tf.word_wrap = True
    style_paragraph(tf.paragraphs[0], "Key Challenges in Deep Networks", HEADER_FONT, 18, ACCENT_GOLD, bold=True, space_after=12)

    p = tf.add_paragraph()
    style_paragraph(p, "• Internal Covariate Shift:", BODY_FONT, 14, TEXT_WHITE, bold=True, space_after=2)
    p = tf.add_paragraph()
    style_paragraph(p, "  - Distributions of layer activations shift continuously as weights update during training, forcing later layers to constantly adapt.", BODY_FONT, 13, TEXT_GRAY, space_after=16)

    p = tf.add_paragraph()
    style_paragraph(p, "• Vanishing & Exploding Gradients:", BODY_FONT, 14, TEXT_WHITE, bold=True, space_after=2)
    p = tf.add_paragraph()
    style_paragraph(p, "  - Error signals shrink to near-zero or expand exponentially as they propagate backward through many layers, halting learning.", BODY_FONT, 13, TEXT_GRAY)


    # ==========================================
    # SLIDE 8: Deep MLP Architecture and Setup
    # ==========================================
    slide = prs.slides.add_slide(slide_layout)
    add_slide_header(slide, "Deep MLP Architecture and Setup")

    # Card 1: Setup details
    add_card(slide, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.0), BORDER_SUBTLE)
    box1 = slide.shapes.add_textbox(Inches(1.05), Inches(1.85), Inches(5.1), Inches(4.5))
    tf = box1.text_frame
    tf.word_wrap = True
    style_paragraph(tf.paragraphs[0], "Dataset & Model Configuration", HEADER_FONT, 18, ACCENT_GOLD, bold=True, space_after=12)

    p = tf.add_paragraph()
    style_paragraph(p, "• Dataset: MNIST Handwritten Digits", BODY_FONT, 14, TEXT_WHITE, bold=True, space_after=2)
    p = tf.add_paragraph()
    style_paragraph(p, "  - 60,000 training images, 10,000 test images (28x28 grayscale digits).", BODY_FONT, 13, TEXT_GRAY, space_after=12)

    p = tf.add_paragraph()
    style_paragraph(p, "• Network Dimensions:", BODY_FONT, 14, TEXT_WHITE, bold=True, space_after=2)
    p = tf.add_paragraph()
    style_paragraph(p, "  - Input Layer: Flattened image vector of size 784.", BODY_FONT, 13, TEXT_GRAY, space_after=4)
    p = tf.add_paragraph()
    style_paragraph(p, "  - Hidden Layers: 5 layers with 512 hidden units each.", BODY_FONT, 13, TEXT_GRAY, space_after=4)
    p = tf.add_paragraph()
    style_paragraph(p, "  - Output Layer: 10 units with Softmax for digit class logits.", BODY_FONT, 13, TEXT_GRAY, space_after=4)
    p = tf.add_paragraph()
    style_paragraph(p, "  - Activation Function: Rectified Linear Unit (ReLU) after each layer.", BODY_FONT, 13, TEXT_GRAY)

    # Card 2: Comparative variables
    add_card(slide, Inches(6.933), Inches(1.6), Inches(5.6), Inches(5.0), CARD_BORDER)
    box2 = slide.shapes.add_textbox(Inches(7.183), Inches(1.85), Inches(5.1), Inches(4.5))
    tf = box2.text_frame
    tf.word_wrap = True
    style_paragraph(tf.paragraphs[0], "Comparative Experimental Setups", HEADER_FONT, 18, ACCENT_GOLD, bold=True, space_after=12)

    p = tf.add_paragraph()
    style_paragraph(p, "• Setup A: Baseline Network (No BN)", BODY_FONT, 14, TEXT_WHITE, bold=True, space_after=2)
    p = tf.add_paragraph()
    style_paragraph(p, "  - Fully Connected (Linear) -> ReLU Activation.", BODY_FONT, 13, TEXT_GRAY, space_after=16)

    p = tf.add_paragraph()
    style_paragraph(p, "• Setup B: Normalized Network (With BN)", BODY_FONT, 14, TEXT_WHITE, bold=True, space_after=2)
    p = tf.add_paragraph()
    style_paragraph(p, "  - Fully Connected -> BatchNorm1d -> ReLU Activation.", BODY_FONT, 13, TEXT_GRAY, space_after=16)

    p = tf.add_paragraph()
    style_paragraph(p, "• Optimization settings:", BODY_FONT, 14, TEXT_WHITE, bold=True, space_after=2)
    p = tf.add_paragraph()
    style_paragraph(p, "  - Trained both setups using standard Mini-batch Stochastic Gradient Descent (SGD) with batch size of 64.", BODY_FONT, 13, TEXT_GRAY)


    # ==========================================
    # SLIDE 9: Experiment 1: Standard Training
    # ==========================================
    slide = prs.slides.add_slide(slide_layout)
    add_slide_header(slide, "Experiment 1: Standard Training")

    # Left Card: Info
    add_card(slide, Inches(0.8), Inches(1.6), Inches(5.2), Inches(5.0), BORDER_SUBTLE)
    box1 = slide.shapes.add_textbox(Inches(1.05), Inches(1.85), Inches(4.7), Inches(4.5))
    tf = box1.text_frame
    tf.word_wrap = True
    style_paragraph(tf.paragraphs[0], "Standard SGD Training", HEADER_FONT, 18, ACCENT_GOLD, bold=True, space_after=12)

    p = tf.add_paragraph()
    style_paragraph(p, "• Hyperparameters:", BODY_FONT, 14, TEXT_WHITE, bold=True, space_after=2)
    p = tf.add_paragraph()
    style_paragraph(p, "  - Learning rate = 0.01, standard initialization, trained for 5 epochs.", BODY_FONT, 13, TEXT_GRAY, space_after=12)

    p = tf.add_paragraph()
    style_paragraph(p, "• Observations:", BODY_FONT, 14, TEXT_WHITE, bold=True, space_after=2)
    p = tf.add_paragraph()
    style_paragraph(p, "  - With Batch Norm: Converges much faster, reaching 97.10% test accuracy.", BODY_FONT, 13, TEXT_GRAY, space_after=4)
    p = tf.add_paragraph()
    style_paragraph(p, "  - No Batch Norm: Converges slowly, reaching 92.83% test accuracy.", BODY_FONT, 13, TEXT_GRAY, space_after=12)

    p = tf.add_paragraph()
    style_paragraph(p, "• Key Takeaway:", BODY_FONT, 14, TEXT_WHITE, bold=True, space_after=2)
    p = tf.add_paragraph()
    style_paragraph(p, "  - Batch Norm acts as an optimization accelerator, helping the model reach higher accuracy in fewer epochs.", BODY_FONT, 13, TEXT_GRAY)

    # Right Card: Image Frame
    add_card(slide, Inches(6.4), Inches(1.6), Inches(6.133), Inches(5.0), CARD_BORDER)
    # Add Image
    img_path = "extracted_plots/task2_exp1_standard.png"
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(6.6), Inches(1.8), Inches(5.733), Inches(4.6))


    # ==========================================
    # SLIDE 10: Experiment 2: Stress-Test with High Learning Rate
    # ==========================================
    slide = prs.slides.add_slide(slide_layout)
    add_slide_header(slide, "Experiment 2: High Learning Rate Stress Test")

    # Left Card: Info
    add_card(slide, Inches(0.8), Inches(1.6), Inches(5.2), Inches(5.0), BORDER_SUBTLE)
    box1 = slide.shapes.add_textbox(Inches(1.05), Inches(1.85), Inches(4.7), Inches(4.5))
    tf = box1.text_frame
    tf.word_wrap = True
    style_paragraph(tf.paragraphs[0], "High Learning Rate Stress Test", HEADER_FONT, 18, ACCENT_GOLD, bold=True, space_after=12)

    p = tf.add_paragraph()
    style_paragraph(p, "• Hyperparameters:", BODY_FONT, 14, TEXT_WHITE, bold=True, space_after=2)
    p = tf.add_paragraph()
    style_paragraph(p, "  - Learning rate = 0.2, standard initialization, trained for 5 epochs.", BODY_FONT, 13, TEXT_GRAY, space_after=12)

    p = tf.add_paragraph()
    style_paragraph(p, "• Observations:", BODY_FONT, 14, TEXT_WHITE, bold=True, space_after=2)
    p = tf.add_paragraph()
    style_paragraph(p, "  - With Batch Norm: Trains stably and converges to 98.24% test accuracy (highest accuracy achieved).", BODY_FONT, 13, TEXT_GRAY, space_after=4)
    p = tf.add_paragraph()
    style_paragraph(p, "  - No Batch Norm: Fails completely to learn. Loss flatlines at 2.3026, test accuracy stays at 11.35% (random guessing).", BODY_FONT, 13, TEXT_GRAY, space_after=12)

    p = tf.add_paragraph()
    style_paragraph(p, "• Key Takeaway:", BODY_FONT, 14, TEXT_WHITE, bold=True, space_after=2)
    p = tf.add_paragraph()
    style_paragraph(p, "  - Batch Norm scale-invariance stabilizes gradient updates, permitting higher learning rates that accelerate training.", BODY_FONT, 13, TEXT_GRAY)

    # Right Card: Image Frame
    add_card(slide, Inches(6.4), Inches(1.6), Inches(6.133), Inches(5.0), CARD_BORDER)
    # Add Image
    img_path = "extracted_plots/task2_exp2_high_lr.png"
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(6.6), Inches(1.8), Inches(5.733), Inches(4.6))


    # ==========================================
    # SLIDE 11: Experiment 3: Stress-Test with Poor Initialization
    # ==========================================
    slide = prs.slides.add_slide(slide_layout)
    add_slide_header(slide, "Experiment 3: Poor Initialization Stress Test")

    # Left Card: Info
    add_card(slide, Inches(0.8), Inches(1.6), Inches(5.2), Inches(5.0), BORDER_SUBTLE)
    box1 = slide.shapes.add_textbox(Inches(1.05), Inches(1.85), Inches(4.7), Inches(4.5))
    tf = box1.text_frame
    tf.word_wrap = True
    style_paragraph(tf.paragraphs[0], "Suboptimal Initialization scale", HEADER_FONT, 18, ACCENT_GOLD, bold=True, space_after=12)

    p = tf.add_paragraph()
    style_paragraph(p, "• Hyperparameters:", BODY_FONT, 14, TEXT_WHITE, bold=True, space_after=2)
    p = tf.add_paragraph()
    style_paragraph(p, "  - Learning rate = 0.01, weights initialized scaled to near-zero (stddev = 0.001), 5 epochs.", BODY_FONT, 13, TEXT_GRAY, space_after=12)

    p = tf.add_paragraph()
    style_paragraph(p, "• Observations:", BODY_FONT, 14, TEXT_WHITE, bold=True, space_after=2)
    p = tf.add_paragraph()
    style_paragraph(p, "  - With Batch Norm: Successfully rescues training, converging to 97.48% test accuracy.", BODY_FONT, 13, TEXT_GRAY, space_after=4)
    p = tf.add_paragraph()
    style_paragraph(p, "  - No Batch Norm: Fails completely. Loss remains stuck at 2.3026, test accuracy stuck at 11.35% (vanishing gradients).", BODY_FONT, 13, TEXT_GRAY, space_after=12)

    p = tf.add_paragraph()
    style_paragraph(p, "• Key Takeaway:", BODY_FONT, 14, TEXT_WHITE, bold=True, space_after=2)
    p = tf.add_paragraph()
    style_paragraph(p, "  - Batch Norm rescales the activations at each layer, maintaining gradient magnitude and preventing vanishing gradients.", BODY_FONT, 13, TEXT_GRAY)

    # Right Card: Image Frame
    add_card(slide, Inches(6.4), Inches(1.6), Inches(6.133), Inches(5.0), CARD_BORDER)
    # Add Image
    img_path = "extracted_plots/task2_exp3_poor_init.png"
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(6.6), Inches(1.8), Inches(5.733), Inches(4.6))


    # ==========================================
    # SLIDE 12: How Batch Normalization Stabilizes Gradients
    # ==========================================
    slide = prs.slides.add_slide(slide_layout)
    add_slide_header(slide, "How Batch Normalization Stabilizes Gradients")

    # Card 1: Math equations
    add_card(slide, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.0), CARD_BORDER)
    box1 = slide.shapes.add_textbox(Inches(1.05), Inches(1.85), Inches(5.1), Inches(4.5))
    tf = box1.text_frame
    tf.word_wrap = True
    style_paragraph(tf.paragraphs[0], "Feedforward Normalization Equations", HEADER_FONT, 18, ACCENT_GOLD, bold=True, space_after=10)

    p = tf.add_paragraph()
    style_paragraph(p, "• Step 1: Mini-batch Mean", BODY_FONT, 13, TEXT_WHITE, bold=True, space_after=2)
    p = tf.add_paragraph()
    style_paragraph(p, "     μ_B = (1/m) * Σᵢ₌₁ᵐ xᵢ", HEADER_FONT, 20, ACCENT_GOLD, bold=True, space_after=8)

    p = tf.add_paragraph()
    style_paragraph(p, "• Step 2: Mini-batch Variance", BODY_FONT, 13, TEXT_WHITE, bold=True, space_after=2)
    p = tf.add_paragraph()
    style_paragraph(p, "     σ_B² = (1/m) * Σᵢ₌₁ᵐ (xᵢ - μ_B)²", HEADER_FONT, 20, ACCENT_GOLD, bold=True, space_after=8)

    p = tf.add_paragraph()
    style_paragraph(p, "• Step 3: Normalize Activation", BODY_FONT, 13, TEXT_WHITE, bold=True, space_after=2)
    p = tf.add_paragraph()
    style_paragraph(p, "     x̂ᵢ = (xᵢ - μ_B) / √(σ_B² + ε)", HEADER_FONT, 20, ACCENT_GOLD, bold=True, space_after=8)

    p = tf.add_paragraph()
    style_paragraph(p, "• Step 4: Scale and Shift (Learnable)", BODY_FONT, 13, TEXT_WHITE, bold=True, space_after=2)
    p = tf.add_paragraph()
    style_paragraph(p, "     yᵢ = γ * x̂ᵢ + β", HEADER_FONT, 20, ACCENT_GOLD, bold=True)

    # Card 2: Why it prevents failures
    add_card(slide, Inches(6.933), Inches(1.6), Inches(5.6), Inches(5.0), BORDER_SUBTLE)
    box2 = slide.shapes.add_textbox(Inches(7.183), Inches(1.85), Inches(5.1), Inches(4.5))
    tf = box2.text_frame
    tf.word_wrap = True
    style_paragraph(tf.paragraphs[0], "Why it Stabilizes Gradient Flow", HEADER_FONT, 18, ACCENT_GOLD, bold=True, space_after=12)

    p = tf.add_paragraph()
    style_paragraph(p, "• Decoupling Weight scale from Activation scale:", BODY_FONT, 14, TEXT_WHITE, bold=True, space_after=2)
    p = tf.add_paragraph()
    style_paragraph(p, "  - Dividing by standard deviation ensures activations never shrink to zero or blow up to infinity, regardless of weight magnitudes.", BODY_FONT, 13, TEXT_GRAY, space_after=12)

    p = tf.add_paragraph()
    style_paragraph(p, "• Preserving Representation Capacity (γ and β):", BODY_FONT, 14, TEXT_WHITE, bold=True, space_after=2)
    p = tf.add_paragraph()
    style_paragraph(p, "  - Learnable scale (γ) and shift (β) parameters let the network restore the original distribution if it is optimal for learning.", BODY_FONT, 13, TEXT_GRAY, space_after=12)

    p = tf.add_paragraph()
    style_paragraph(p, "• Smoother Loss Landscape:", BODY_FONT, 14, TEXT_WHITE, bold=True, space_after=2)
    p = tf.add_paragraph()
    style_paragraph(p, "  - Standardizing intermediate layers prevents extreme gradient fluctuations, leading to a much smoother optimization path.", BODY_FONT, 13, TEXT_GRAY)


    # ==========================================
    # SLIDE 13: Summary of Empirical Findings (TABLE!)
    # ==========================================
    slide = prs.slides.add_slide(slide_layout)
    add_slide_header(slide, "Summary of Empirical Findings")

    # Native PPTX Table in the middle
    rows, cols = 4, 4
    left = Inches(1.5)
    top = Inches(1.8)
    width = Inches(10.333)
    height = Inches(3.2)

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    # Column Widths
    table.columns[0].width = Inches(3.5)
    table.columns[1].width = Inches(1.8)
    table.columns[2].width = Inches(2.5)
    table.columns[3].width = Inches(2.533)

    # Table Contents
    data = [
        ["Experiment", "Metric", "No Batch Norm (Baseline)", "With Batch Norm"],
        ["Standard Training (LR=0.01)", "Test Accuracy", "92.83%", "97.10% (Accelerated)"],
        ["High Learning Rate Stress Test (LR=0.2)", "Test Accuracy", "11.35% (Diverged)", "98.24% (Stable)"],
        ["Suboptimal Initialization (σ=0.001)", "Test Accuracy", "11.35% (Vanished)", "97.48% (Rescued)"]
    ]

    for r_idx, row in enumerate(data):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = val
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            # Format paragraph
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            
            # Formatting styles
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER_BG
                style_paragraph(p, val, HEADER_FONT, 15, TEXT_WHITE, bold=True, space_after=0)
            else:
                cell.fill.solid()
                if r_idx % 2 == 1:
                    cell.fill.fore_color.rgb = RGBColor(26, 32, 44)
                else:
                    cell.fill.fore_color.rgb = RGBColor(16, 22, 34)
                
                # Check if BN column for bold highlight
                is_bn_col = (c_idx == 3)
                txt_color = ACCENT_GOLD if (is_bn_col and r_idx > 0) else TEXT_GRAY
                style_paragraph(p, val, BODY_FONT, 14, txt_color, bold=is_bn_col, space_after=0)

    # Add Takeaway at the bottom
    takeaway_box = slide.shapes.add_textbox(Inches(1.5), Inches(5.3), Inches(10.333), Inches(1.4))
    tf = takeaway_box.text_frame
    tf.word_wrap = True
    style_paragraph(tf.paragraphs[0], "Key Empirical Findings", HEADER_FONT, 16, ACCENT_GOLD, bold=True, space_after=4)
    p = tf.add_paragraph()
    style_paragraph(p, "• Batch Normalization consistently out-performs the baseline in standard settings and successfully rescues the network from extreme training failures (learning rate explosion and weight initialization vanishing).", BODY_FONT, 13, TEXT_GRAY)


    # ==========================================
    # SLIDE 14: Connections to Deep Learning Theory (Weeks 1-7)
    # ==========================================
    slide = prs.slides.add_slide(slide_layout)
    add_slide_header(slide, "Connections to Deep Learning Theory")

    # Card 1: Universal Approx & Init
    add_card(slide, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.0), BORDER_SUBTLE)
    box1 = slide.shapes.add_textbox(Inches(1.05), Inches(1.85), Inches(5.1), Inches(4.5))
    tf = box1.text_frame
    tf.word_wrap = True
    style_paragraph(tf.paragraphs[0], "Universal Approximation & Initialization", HEADER_FONT, 18, ACCENT_GOLD, bold=True, space_after=12)

    p = tf.add_paragraph()
    style_paragraph(p, "• Universal Approximation Theorem:", BODY_FONT, 14, TEXT_WHITE, bold=True, space_after=2)
    p = tf.add_paragraph()
    style_paragraph(p, "  - Standard deep feedforward networks and Recurrent structures can approximate arbitrary complex non-linear mappings. Character-level LSTMs utilize this representation capacity to discover language syntax dynamically from character context.", BODY_FONT, 12, TEXT_GRAY, space_after=16)

    p = tf.add_paragraph()
    style_paragraph(p, "• Weight Initialization (Xavier / He):", BODY_FONT, 14, TEXT_WHITE, bold=True, space_after=2)
    p = tf.add_paragraph()
    style_paragraph(p, "  - Proper initialization (e.g. He initialization for ReLU) is critical to prevent activation variance from vanishing or exploding in deep layers. Batch Norm relaxes this structural requirement by dynamically rescaling the variance at every layer.", BODY_FONT, 12, TEXT_GRAY)

    # Card 2: Optimization landscapes
    add_card(slide, Inches(6.933), Inches(1.6), Inches(5.6), Inches(5.0), CARD_BORDER)
    box2 = slide.shapes.add_textbox(Inches(7.183), Inches(1.85), Inches(5.1), Inches(4.5))
    tf = box2.text_frame
    tf.word_wrap = True
    style_paragraph(tf.paragraphs[0], "Optimization Dynamics & Regularization", HEADER_FONT, 18, ACCENT_GOLD, bold=True, space_after=12)

    p = tf.add_paragraph()
    style_paragraph(p, "• Optimization Landscapes:", BODY_FONT, 14, TEXT_WHITE, bold=True, space_after=2)
    p = tf.add_paragraph()
    style_paragraph(p, "  - High learning rates cause gradient updates to overshoot local minima and diverge. Batch Norm smoothens the optimization landscape, preventing gradients from exploding and making optimization highly stable.", BODY_FONT, 12, TEXT_GRAY, space_after=16)

    p = tf.add_paragraph()
    style_paragraph(p, "• Implicit Regularization:", BODY_FONT, 14, TEXT_WHITE, bold=True, space_after=2)
    p = tf.add_paragraph()
    style_paragraph(p, "  - The calculation of mean and variance over finite mini-batches introduces moderate stochastic noise into layer activations. This noise acts as a regularizer, reducing overfitting similar to Dropout.", BODY_FONT, 12, TEXT_GRAY)


    # ==========================================
    # SLIDE 15: Conclusion and Q&A
    # ==========================================
    slide = prs.slides.add_slide(slide_layout)
    add_slide_header(slide, "Conclusion and Q&A")

    # Card 1: Summary points
    add_card(slide, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.0), BORDER_SUBTLE)
    box1 = slide.shapes.add_textbox(Inches(1.05), Inches(1.85), Inches(5.1), Inches(4.5))
    tf = box1.text_frame
    tf.word_wrap = True
    style_paragraph(tf.paragraphs[0], "Key Summary Takeaways", HEADER_FONT, 18, ACCENT_GOLD, bold=True, space_after=12)

    p = tf.add_paragraph()
    style_paragraph(p, "• Sequence Modeling with LSTMs:", BODY_FONT, 14, TEXT_WHITE, bold=True, space_after=2)
    p = tf.add_paragraph()
    style_paragraph(p, "  - Stacked LSTMs can successfully learn the structural, orthographical, and grammatical rules of natural text entirely from character sequences without pre-defined vocabularies.", BODY_FONT, 13, TEXT_GRAY, space_after=16)

    p = tf.add_paragraph()
    style_paragraph(p, "• Activation Normalization with Batch Norm:", BODY_FONT, 14, TEXT_WHITE, bold=True, space_after=2)
    p = tf.add_paragraph()
    style_paragraph(p, "  - Batch Normalization is a critical architectural element that stabilizes training, accelerates convergence, and provides extreme robustness to learning rates and initialization scales.", BODY_FONT, 13, TEXT_GRAY)

    # Card 2: Thank you panel
    add_card(slide, Inches(6.933), Inches(1.6), Inches(5.6), Inches(5.0), CARD_BORDER)
    box2 = slide.shapes.add_textbox(Inches(7.183), Inches(1.85), Inches(5.1), Inches(4.5))
    tf = box2.text_frame
    tf.word_wrap = True
    style_paragraph(tf.paragraphs[0], "Thank you! Questions?", HEADER_FONT, 22, ACCENT_GOLD, bold=True, space_after=16)

    p = tf.add_paragraph()
    style_paragraph(p, "Presenters Panel:", BODY_FONT, 14, TEXT_WHITE, bold=True, space_after=8)

    presenters = [
        "MD Kaif Molla",
        "Arhan Das",
        "Sumit Akhuli",
        "Sarthak Mishra",
        "Netram Faran",
        "Aditya Prasad"
    ]
    for pres in presenters:
        p = tf.add_paragraph()
        style_paragraph(p, "  • " + pres, BODY_FONT, 13, TEXT_GRAY, space_after=4)

    # Save presentation
    prs.save("presentation.pptx")
    print("Presentation successfully saved to presentation.pptx!")

if __name__ == "__main__":
    create_presentation()
